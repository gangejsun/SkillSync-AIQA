#!/usr/bin/env python3
"""
PPT Generator - Automatic PowerPoint Creation
Creates professional presentations from structured text input
"""

import sys
import argparse
import re
from pathlib import Path
from typing import List, Dict, Tuple
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor


# Color Themes
COLOR_THEMES = {
    'blue': {
        'primary': (46, 80, 144),      # #2E5090
        'secondary': (74, 144, 226),   # #4A90E2
        'accent': (124, 185, 232),     # #7CB9E8
        'background': (240, 244, 248), # #F0F4F8
        'text': (33, 33, 33),          # #212121
    },
    'green': {
        'primary': (45, 95, 63),       # #2D5F3F
        'secondary': (76, 175, 80),    # #4CAF50
        'accent': (129, 199, 132),     # #81C784
        'background': (241, 248, 244), # #F1F8F4
        'text': (33, 33, 33),
    },
    'purple': {
        'primary': (94, 53, 177),      # #5E35B1
        'secondary': (126, 87, 194),   # #7E57C2
        'accent': (179, 157, 219),     # #B39DDB
        'background': (243, 240, 248), # #F3F0F8
        'text': (33, 33, 33),
    },
    'orange': {
        'primary': (216, 67, 21),      # #D84315
        'secondary': (255, 112, 67),   # #FF7043
        'accent': (255, 171, 145),     # #FFAB91
        'background': (255, 243, 224), # #FFF3E0
        'text': (33, 33, 33),
    }
}


def parse_slides(content: str) -> List[Dict]:
    """Parse structured text into slide data"""
    slides = []
    current_slide = None

    lines = content.strip().split('\n')

    for line in lines:
        line = line.rstrip()

        # Slide header: "Slide X: [타입]"
        slide_match = re.match(r'^Slide\s+\d+:\s*\[(.+?)\]', line, re.IGNORECASE)
        if slide_match:
            if current_slide:
                slides.append(current_slide)

            slide_type = slide_match.group(1).strip()
            current_slide = {
                'type': slide_type,
                'title': '',
                'content': [],
                'left_content': [],
                'right_content': [],
                'current_section': 'content'
            }
            continue

        if not current_slide:
            continue

        # Title line
        if line.startswith('제목:') or line.startswith('title:'):
            current_slide['title'] = line.split(':', 1)[1].strip()
            continue

        # Subtitle line (for title slides)
        if line.startswith('부제:') or line.startswith('subtitle:'):
            current_slide['subtitle'] = line.split(':', 1)[1].strip()
            continue

        # Section markers for 2-column layout
        if line.strip() == '[왼쪽]' or line.strip() == '[left]':
            current_slide['current_section'] = 'left'
            continue

        if line.strip() == '[오른쪽]' or line.strip() == '[right]':
            current_slide['current_section'] = 'right'
            continue

        # Content lines (bullets, text, etc.)
        if line.strip():
            if current_slide['current_section'] == 'left':
                current_slide['left_content'].append(line)
            elif current_slide['current_section'] == 'right':
                current_slide['right_content'].append(line)
            else:
                current_slide['content'].append(line)

    # Add last slide
    if current_slide:
        slides.append(current_slide)

    return slides


def create_title_slide(prs, slide_data: Dict, theme: Dict):
    """Create title slide"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*theme['background'])

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(1), Inches(2.5), Inches(8), Inches(1.5)
    )
    title_frame = title_box.text_frame
    title_frame.text = slide_data['title']
    title_frame.paragraphs[0].font.size = Pt(54)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(*theme['primary'])
    title_frame.paragraphs[0].font.name = '맑은 고딕'
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Subtitle
    if 'subtitle' in slide_data:
        subtitle_box = slide.shapes.add_textbox(
            Inches(1), Inches(4.2), Inches(8), Inches(1)
        )
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = slide_data['subtitle']
        subtitle_frame.paragraphs[0].font.size = Pt(28)
        subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(*theme['secondary'])
        subtitle_frame.paragraphs[0].font.name = '맑은 고딕'
        subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Accent line
    line = slide.shapes.add_shape(
        1,  # Line shape
        Inches(3.5), Inches(4.0), Inches(3), Inches(0)
    )
    line.line.color.rgb = RGBColor(*theme['accent'])
    line.line.width = Pt(3)


def create_content_slide(prs, slide_data: Dict, theme: Dict):
    """Create standard content slide with bullets"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*theme['background'])

    # Title
    if slide_data['title']:
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.4), Inches(9), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.text = slide_data['title']
        title_frame.paragraphs[0].font.size = Pt(44)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(*theme['primary'])
        title_frame.paragraphs[0].font.name = '맑은 고딕'

    # Content
    content_box = slide.shapes.add_textbox(
        Inches(0.8), Inches(1.5), Inches(8.4), Inches(4.5)
    )
    text_frame = content_box.text_frame
    text_frame.word_wrap = True

    for line in slide_data['content']:
        # Determine indent level
        indent_level = 0
        cleaned_line = line.lstrip()

        if cleaned_line.startswith('- '):
            indent_level = 1
            cleaned_line = cleaned_line[2:]
        elif cleaned_line.startswith('• '):
            indent_level = 0
            cleaned_line = cleaned_line[2:]

        # Add paragraph
        p = text_frame.add_paragraph()
        p.text = cleaned_line
        p.level = indent_level
        p.font.size = Pt(24 if indent_level == 0 else 20)
        p.font.color.rgb = RGBColor(*theme['text'])
        p.font.name = '맑은 고딕'
        p.space_before = Pt(6)


def create_two_column_slide(prs, slide_data: Dict, theme: Dict):
    """Create two-column comparison slide"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*theme['background'])

    # Title
    if slide_data['title']:
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.4), Inches(9), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.text = slide_data['title']
        title_frame.paragraphs[0].font.size = Pt(44)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(*theme['primary'])
        title_frame.paragraphs[0].font.name = '맑은 고딕'

    # Left column
    left_box = slide.shapes.add_textbox(
        Inches(0.8), Inches(1.5), Inches(4), Inches(4.5)
    )
    left_frame = left_box.text_frame
    left_frame.word_wrap = True

    for line in slide_data['left_content']:
        p = left_frame.add_paragraph()
        p.text = line.lstrip('• -').strip()
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(*theme['text'])
        p.font.name = '맑은 고딕'
        p.space_before = Pt(6)

    # Right column
    right_box = slide.shapes.add_textbox(
        Inches(5.2), Inches(1.5), Inches(4), Inches(4.5)
    )
    right_frame = right_box.text_frame
    right_frame.word_wrap = True

    for line in slide_data['right_content']:
        p = right_frame.add_paragraph()
        p.text = line.lstrip('• -').strip()
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(*theme['text'])
        p.font.name = '맑은 고딕'
        p.space_before = Pt(6)

    # Divider line
    line = slide.shapes.add_shape(
        1,  # Line shape
        Inches(5), Inches(1.5), Inches(0), Inches(4.5)
    )
    line.line.color.rgb = RGBColor(*theme['accent'])
    line.line.width = Pt(2)


def create_quote_slide(prs, slide_data: Dict, theme: Dict):
    """Create quote/highlight slide"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*theme['primary'])

    # Quote text
    quote_text = '\n'.join(slide_data['content'])

    quote_box = slide.shapes.add_textbox(
        Inches(1.5), Inches(2), Inches(7), Inches(3)
    )
    text_frame = quote_box.text_frame
    text_frame.text = quote_text
    text_frame.paragraphs[0].font.size = Pt(36)
    text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    text_frame.paragraphs[0].font.name = '맑은 고딕'
    text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE


def generate_presentation(slides_data: List[Dict], output_path: str, color: str = 'blue', style: str = 'professional'):
    """Generate PowerPoint presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    theme = COLOR_THEMES.get(color, COLOR_THEMES['blue'])

    for slide_data in slides_data:
        slide_type = slide_data['type'].lower()

        if '제목' in slide_type or 'title' in slide_type:
            create_title_slide(prs, slide_data, theme)
        elif '2단' in slide_type or 'column' in slide_type or '비교' in slide_type:
            create_two_column_slide(prs, slide_data, theme)
        elif '인용' in slide_type or 'quote' in slide_type:
            create_quote_slide(prs, slide_data, theme)
        else:
            # Default: content slide
            create_content_slide(prs, slide_data, theme)

    # Save presentation
    prs.save(output_path)
    print(f"✅ Presentation saved: {output_path}")
    print(f"   Slides: {len(slides_data)}")
    print(f"   Color: {color}")
    print(f"   Style: {style}")


def main():
    parser = argparse.ArgumentParser(description='Generate PowerPoint presentations')
    parser.add_argument('--output', required=True, help='Output PPTX file path')
    parser.add_argument('--style', default='professional', choices=['professional', 'minimal', 'creative'],
                        help='Presentation style')
    parser.add_argument('--color', default='blue', choices=['blue', 'green', 'purple', 'orange'],
                        help='Color theme')

    args = parser.parse_args()

    # Read content from stdin
    content = sys.stdin.read()

    if not content.strip():
        print("Error: No content provided via stdin", file=sys.stderr)
        sys.exit(1)

    # Parse slides
    slides_data = parse_slides(content)

    if not slides_data:
        print("Error: No valid slides found in input", file=sys.stderr)
        sys.exit(1)

    # Create output directory if needed
    output_path = Path(args.output)
    output_path.parent.mkdir(parents=True, exist_ok=True)

    # Generate presentation
    generate_presentation(slides_data, str(output_path), args.color, args.style)


if __name__ == '__main__':
    main()
